"""Generate Enterprise Presentation PPTX for Elevate HR (Team 12)."""
import os
import sys

# Ensure site-packages are loaded
sys.path.insert(0, "/Users/zuhaibp/Documents/Project_elevate_team_12/.venv/lib/python3.11/site-packages")

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_PPTX = "/Users/zuhaibp/Documents/Project_elevate_team_12/Elevate_HR_Team_12_Autonomous_MAS_Presentation.pptx"

# Modern Color Palette (Google / Enterprise Tech Theme)
C_DARK_BG = RGBColor(15, 23, 42)       # Slate 900
C_LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
C_WHITE = RGBColor(255, 255, 255)
C_CARD_BG = RGBColor(255, 255, 255)
C_CARD_BORDER = RGBColor(226, 232, 240)
C_PRIMARY = RGBColor(26, 115, 232)     # Google Blue #1A73E8
C_ACCENT_RED = RGBColor(234, 67, 53)   # Google Red #EA4335
C_ACCENT_YELLOW = RGBColor(251, 188, 5)# Google Yellow #FBBC05
C_ACCENT_GREEN = RGBColor(52, 168, 83) # Google Green #34A853
C_TEXT_MAIN = RGBColor(30, 41, 59)     # Slate 800
C_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
C_PRIMARY_LIGHT = RGBColor(232, 240, 254)

# Available Image Paths
IMG_ARCH = "/Users/zuhaibp/.gemini/jetski/brain/922a6cbf-9997-4285-b616-0805f372d1b5/system_architecture_1786964663566.jpg"
IMG_FLOW = "/Users/zuhaibp/.gemini/jetski/brain/922a6cbf-9997-4285-b616-0805f372d1b5/flow_diagram_1786964693259.jpg"
IMG_PORTAL = "/Users/zuhaibp/.gemini/jetski/brain/922a6cbf-9997-4285-b616-0805f372d1b5/.user_uploaded/media_1787200239248.png"
IMG_UI_APP = "/Users/zuhaibp/.gemini/jetski/brain/922a6cbf-9997-4285-b616-0805f372d1b5/.user_uploaded/media_1787200246731.png"


def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen Layout (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Helper: Set slide background color
    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Add Standard Header
    def add_header(slide, category, title, dark=False):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Category / Eyebrow
        p_cat = tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_PRIMARY if not dark else C_ACCENT_YELLOW
        p_cat.space_after = Pt(2)

        # Title
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = C_TEXT_MAIN if not dark else C_WHITE

    # Helper: Add Styled Card / Box
    def add_card(slide, left, top, width, height, bg_color=C_WHITE, border_color=C_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape


    # =========================================================================
    # SLIDE 8.5: The Query Flow & State Management
    # =========================================================================
    s8_5 = prs.slides.add_slide(blank_layout)
    set_bg(s8_5, C_LIGHT_BG)
    add_header(s8_5, "SYSTEM ARCHITECTURE", "End-to-End Query Flow & State Management")

    # Flow Box
    f_box = add_card(s8_5, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    f_tf = f_box.text_frame
    f_tf.word_wrap = True
    f_tf.margin_left = f_tf.margin_top = f_tf.margin_right = f_tf.margin_bottom = Inches(0.3)
    
    p = f_tf.paragraphs[0]
    p.text = "How a Request Moves Through the System:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    
    flow_steps = [
        ("1. UI & Auth Injection", "Browser generates chat_session_id. UI injects X-Employee-ID and X-MCP-Token into FastAPI REST headers."),
        ("2. Backend Thread Context", "FastAPI (ui/server.py) maps session ID to ADK memory. Token is injected securely into ACTIVE_MCP_TOKEN_CV ContextVar."),
        ("3. Orchestrator Routing", "Orchestrator receives prompt. Delegates to Sub-Agents (HCM / ITSM / Policy) based on domain context."),
        ("4. Autonomous Multi-Step Execution", "If bulk action requested (e.g., 'Close all tickets'), agent autonomously calls 'Read' tools (list_tickets) first, parses IDs, and loops 'Write' tools (update_ticket_status) without acting like a dumb chatbot."),
        ("5. State Management & Real-Time Sync", "Agent memory stores conversation. SaaS backends store facts. Agent ALWAYS executes live reads to prevent hallucinating stale balances. UI sidebar autonomously fetches real-time updates post-action.")
    ]
    
    for title, desc in flow_steps:
        pt = f_tf.add_paragraph()
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = C_TEXT_MAIN
        pt.space_before = Pt(12)
        
        pd = f_tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(14)
        pd.font.color.rgb = C_TEXT_MUTED
        pd.level = 1

    # =========================================================================
    # SLIDE 8.6: Strict Policy Gatekeeper & Anti-Downgrade
    # =========================================================================
    s8_6 = prs.slides.add_slide(blank_layout)
    set_bg(s8_6, C_DARK_BG)
    add_header(s8_6, "COMPLIANCE", "The Pre-Action Policy Gatekeeper", dark=True)

    g_box = s8_6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    g_tf = g_box.text_frame
    g_tf.word_wrap = True
    
    gp = g_tf.paragraphs[0]
    gp.text = "Strict Two-Step Escalation & No Silent Downgrades"
    gp.font.size = Pt(20)
    gp.font.bold = True
    gp.font.color.rgb = C_ACCENT_YELLOW
    
    guard_steps = [
        ("The Problem:", "Helpful LLMs often silently modify non-compliant requests (e.g. downgrading 'Priority 1' to 'Priority 3') just to force the action through. Or they directly create Escalation Tickets instead of refusing the user."),
        ("The Gatekeeper Solution:", "Before ANY write-tool execution, the agent evaluates policy."),
        ("Step 1: Hard Stop & Deny", "Agent halts entirely. Outputs a prominent ⚠️ Policy Non-Compliance Warning. It is explicitly forbidden from secretly downgrading parameters."),
        ("Step 2: Multi-Part Query Defense", "Agent is forced to answer all secondary questions (e.g. $500 equipment allowances) in the same denial response."),
        ("Step 3: Explicit Escalation", "Agent only verbally suggests escalation. It NEVER executes 'escalate_to_human_hr' unless the user replies back in the next turn confirming it.")
    ]
    
    for title, desc in guard_steps:
        pt = g_tf.add_paragraph()
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = C_WHITE
        pt.space_before = Pt(16)
        
        pd = g_tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(14)
        pd.font.color.rgb = C_LIGHT_BG
        pd.level = 1


    # =========================================================================
    # SLIDE 1: Title Slide (Dark Theme)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, C_DARK_BG)

    # Accent Glow Element
    glow = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = C_PRIMARY
    glow.line.fill.background()

    # Title & Subtitle Box
    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.2), Inches(4.5))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "GOOGLE ENTERPRISE AI HACKATHON • TEAM 12"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = C_ACCENT_YELLOW
    p0.space_after = Pt(12)

    p1 = tf1.add_paragraph()
    p1.text = "Elevate HR"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Autonomous Multi-Agent HR Orchestrator"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = C_PRIMARY
    p2.space_after = Pt(18)

    p3 = tf1.add_paragraph()
    p3.text = "Next-Generation Enterprise Self-Service powered by Google ADK, FastMCP, Model Armor & Open Knowledge Federation."
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(24)

    p4 = tf1.add_paragraph()
    p4.text = "Architectural Blueprint • Decision Choices • Security Guardrails • Live System Verification"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = C_ACCENT_GREEN

    # =========================================================================
    # SLIDE 2: Executive Summary & Enterprise Problem Statement
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, C_LIGHT_BG)
    add_header(s2, "Executive Overview", "Transforming Enterprise HR Friction into Autonomous Resolution")

    cards_data = [
        ("The Enterprise Challenge", C_ACCENT_RED, [
            "Fragmented Silos: Employees navigate separate portals for Policies (PDFs), HCM (Workday/WorkWeek), and ITSM (ServiceNow/ServiceImmediately).",
            "Context Loss: Compound workflows (e.g. sick leave -> PTO booking -> email delegation ticket) require 3 manual hops.",
            "Compliance Vulnerabilities: Employees misclassify expenses or take unapproved leaves due to complex statutory policies."
        ]),
        ("The Agentic Solution", C_PRIMARY, [
            "Hierarchical Multi-Agent Topology: Built on Google Agent Development Kit (ADK) with root orchestrator & 3 domain specialists.",
            "Open Knowledge Federation (OKF): Sub-millisecond policy retrieval with strict URI citations & zero-hallucination guarantees.",
            "FastMCP Integration: Model Context Protocol over Streamable HTTP for live bi-directional HCM and ITSM transactions."
        ]),
        ("Measurable Impact", C_ACCENT_GREEN, [
            "85%+ Deflection of routine Tier-1 HR and IT helpdesk queries.",
            "<10s SLA across multi-turn compound cross-system workflows.",
            "100% Policy Compliance with hard programmatic and prompt invariants ('No means No').",
            "99.34% Score on comprehensive 33-case benchmark evaluation suite."
        ])
    ]

    left_pos = Inches(0.8)
    for title, accent, bullets in cards_data:
        add_card(s2, left_pos, Inches(1.6), Inches(3.64), Inches(5.2))
        
        tb = s2.shapes.add_textbox(left_pos + Inches(0.25), Inches(1.8), Inches(3.14), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = accent
        p.space_after = Pt(14)

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = f"• {b}"
            pb.font.size = Pt(12)
            pb.font.color.rgb = C_TEXT_MAIN
            pb.space_after = Pt(10)

        left_pos += Inches(3.95)

    # =========================================================================
    # SLIDE 3: System Architecture & Topology
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, C_LIGHT_BG)
    add_header(s3, "Architecture Topology", "Decoupled Hierarchical Multi-Agent Architecture")

    if os.path.exists(IMG_ARCH):
        s3.shapes.add_picture(IMG_ARCH, Inches(0.8), Inches(1.6), width=Inches(6.2))

    right_left = Inches(7.3)
    arch_pillars = [
        ("Central ADK Root Orchestrator (`hr_orchestrator`)", [
            "Acts as the intelligent coordinator and single conversational interface.",
            "Performs intent classification, context aggregation, and dynamic delegation.",
            "Synthesizes multi-agent outputs into unified executive responses."
        ]),
        ("Policy Specialist (`policy_specialist`)", [
            "Open Knowledge Format (OKF) with thread-safe atomic cache hot-reloading.",
            "Extracts statutory Singapore MOM rules, leaves, benefits, and code of conduct.",
            "Enforces strict URI citations (`policy://...`) and zero hallucination."
        ]),
        ("HCM Specialist (`hcm_specialist`) & ITSM Specialist (`itsm_specialist`)", [
            "WorkWeek FastMCP: Checks balances, submits time-off, updates profiles.",
            "ServiceImmediately FastMCP: Ticket creation, comment activity, Tier-2 HITL escalation."
        ])
    ]

    top_pos = Inches(1.6)
    for p_title, p_bullets in arch_pillars:
        add_card(s3, right_left, top_pos, Inches(5.2), Inches(1.65))
        tb = s3.shapes.add_textbox(right_left + Inches(0.2), top_pos + Inches(0.1), Inches(4.8), Inches(1.45))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.space_after = Pt(4)

        for b in p_bullets:
            pb = tf.add_paragraph()
            pb.text = f"• {b}"
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = C_TEXT_MAIN
            pb.space_after = Pt(2)

        top_pos += Inches(1.8)

    # =========================================================================
    # SLIDE 4: Decision Choices — Why This and Not That?
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, C_LIGHT_BG)
    add_header(s4, "Engineering Rationale", "Key Architectural Decisions: Why This and Not That?")

    decisions = [
        ("Architecture", "Google ADK Multi-Agent Hierarchy", "Monolithic Prompt / Flat LangChain", [
            "Domain tool isolation prevents tool hallucination and prompt pollution.",
            "Lightweight ADK Runner with standard SessionService for robust state management."
        ]),
        ("Knowledge Engine", "Open Knowledge Federation (OKF)", "Vector DB Embeddings / Heavy RAG", [
            "Atomic double-buffered cache with mtime auto-invalidation (sub-ms lookup).",
            "Eliminates chunking loss, stale embeddings, and ungrounded hallucinations."
        ]),
        ("Tool Integration", "FastMCP over Streamable HTTP", "Custom Proprietary REST Endpoints", [
            "Standardized Model Context Protocol (MCP) dynamic tool schema discovery.",
            "Universal deployment token binding with fallback resilience."
        ]),
        ("Security & DLP", "Model Armor + Regex Sanitizer", "Post-hoc LLM Self-Moderation", [
            "Zero SPII leakage (NRIC/Credit Cards masked BEFORE reaching LLM).",
            "Immutable policy invariants: 'No Means No' even when users demand overrides."
        ])
    ]

    positions = [
        (Inches(0.8), Inches(1.6)),
        (Inches(6.8), Inches(1.6)),
        (Inches(0.8), Inches(4.35)),
        (Inches(6.8), Inches(4.35))
    ]

    card_w = Inches(5.6)
    card_h = Inches(2.55)

    for idx, (cat, chosen, alt, reasons) in enumerate(decisions):
        pos_l, pos_t = positions[idx]
        add_card(s4, pos_l, pos_t, card_w, card_h)

        tb = s4.shapes.add_textbox(pos_l + Inches(0.2), pos_t + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = f"{cat.upper()} DECISION"
        p0.font.size = Pt(10.5)
        p0.font.bold = True
        p0.font.color.rgb = C_PRIMARY
        p0.space_after = Pt(3)

        p1 = tf.add_paragraph()
        p1.text = f"✅ Chosen: {chosen}"
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = C_ACCENT_GREEN

        p2 = tf.add_paragraph()
        p2.text = f"❌ Rejected: {alt}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_after = Pt(6)

        for r in reasons:
            pr = tf.add_paragraph()
            pr.text = f"• {r}"
            pr.font.size = Pt(10)
            pr.font.color.rgb = C_TEXT_MAIN
            pr.space_after = Pt(2)

    # =========================================================================
    # SLIDE 5: Autonomous Multi-System Delegation Flow (Diagram)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, C_LIGHT_BG)
    add_header(s5, "Autonomous Execution Flow", "Compound Multi-Hop Workflow Delegation")

    if os.path.exists(IMG_FLOW):
        s5.shapes.add_picture(IMG_FLOW, Inches(0.8), Inches(1.6), width=Inches(6.2))

    top_w = Inches(1.6)
    workflow_steps = [
        ("Step 1: Ingestion & In-Flight Masking", "User prompt is intercepted by Model Armor regex filter; NRIC, cards, and credentials are redacted prior to LLM submission."),
        ("Step 2: Orchestrator Intent Decomposition", "Root agent analyzes query (e.g. 'I need 2 days leave and a ticket to forward emails') and decomposes into parallel/sequential sub-tasks."),
        ("Step 3: Specialist Tool Invocation", "Policy Specialist verifies notice rules -> HCM Specialist books time off -> ITSM Specialist creates email routing incident."),
        ("Step 4: Synthesis, Citation & UI Stream", "Results are collated with ticket IDs, balance confirmations, and bottom policy sources, streamed back to user in sub-second response.")
    ]

    for s_title, s_desc in workflow_steps:
        add_card(s5, Inches(7.3), top_w, Inches(5.2), Inches(1.2))
        tb = s5.shapes.add_textbox(Inches(7.5), top_w + Inches(0.1), Inches(4.8), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = s_title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = s_desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = C_TEXT_MAIN

        top_w += Inches(1.35)

    # =========================================================================
    # SLIDE 6: Live Web Application & Workspace UI (Screenshot)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, C_LIGHT_BG)
    add_header(s6, "Modern User Experience", "Google Aura 3-Column Workspace UI")

    if os.path.exists(IMG_UI_APP):
        s6.shapes.add_picture(IMG_UI_APP, Inches(0.8), Inches(1.6), width=Inches(6.8))

    add_card(s6, Inches(7.9), Inches(1.6), Inches(4.6), Inches(5.2))
    tb_ui = s6.shapes.add_textbox(Inches(8.15), Inches(1.85), Inches(4.1), Inches(4.7))
    tf_ui = tb_ui.text_frame
    tf_ui.word_wrap = True

    p = tf_ui.paragraphs[0]
    p.text = "Key UI Capabilities"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.space_after = Pt(12)

    ui_points = [
        ("3-Column Layout", "Left: Chat Session History\nCenter: Active Conversation Stream\nRight: Real-time 'My Hub' drawer"),
        ("Live 'My Hub' Sidebar", "Dynamic leave balance bars and recent ITSM ticket tracker decoded live from FastMCP."),
        ("Multi-User Persona Switcher", "Seamless demonstration login across EMP-380 (Zuhaib), EMP-102 (Sarah), EMP-001 (Alex)."),
        ("Execution Trace Viewer", "Collapsible tool-call trace revealing exact JSON-RPC payloads and ADK sub-agent handoffs."),
        ("One-Click Quick Actions", "Direct chip buttons for balance checks, ticket summaries, and policy FAQs.")
    ]

    for u_title, u_desc in ui_points:
        pu = tf_ui.add_paragraph()
        pu.text = f"✨ {u_title}"
        pu.font.size = Pt(12)
        pu.font.bold = True
        pu.font.color.rgb = C_TEXT_MAIN

        pd = tf_ui.add_paragraph()
        pd.text = u_desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = C_TEXT_MUTED
        pd.space_after = Pt(8)

    # =========================================================================
    # SLIDE 7: Live SaaS Integration (ServiceImmediately & WorkWeek)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, C_LIGHT_BG)
    add_header(s7, "Live Enterprise Integrations", "Real-Time FastMCP SaaS Execution & Attribution")

    if os.path.exists(IMG_PORTAL):
        s7.shapes.add_picture(IMG_PORTAL, Inches(0.8), Inches(1.6), width=Inches(7.2))

    add_card(s7, Inches(8.3), Inches(1.6), Inches(4.2), Inches(5.2))
    tb_saas = s7.shapes.add_textbox(Inches(8.55), Inches(1.85), Inches(3.7), Inches(4.7))
    tf_saas = tb_saas.text_frame
    tf_saas.word_wrap = True

    p = tf_saas.paragraphs[0]
    p.text = "Verified SaaS Workflows"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_GREEN
    p.space_after = Pt(12)

    saas_points = [
        ("Direct Portal Creation", "Tickets (INC0002836, INC0002835, etc.) appear instantly on the live SaaS portal with zero human middleware."),
        ("Accurate Attribution", "Initial dispatch comments automatically link the employee's ID (e.g. 'Updated by: 380') in the audit trail."),
        ("Real-Time Balances", "Leave deductions in WorkWeek immediately update accrued vacation/sick balances in the UI drawer."),
        ("Resilient Fallback", "High-fidelity mock execution gracefully protects user experience during sandbox credential rotations.")
    ]

    for st, sd in saas_points:
        pu = tf_saas.add_paragraph()
        pu.text = f"🔗 {st}"
        pu.font.size = Pt(11.5)
        pu.font.bold = True
        pu.font.color.rgb = C_TEXT_MAIN

        pd = tf_saas.add_paragraph()
        pd.text = sd
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = C_TEXT_MUTED
        pd.space_after = Pt(8)

    # =========================================================================
    # SLIDE 8: Strict Compliance Guardrails ("No Means No")
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, C_LIGHT_BG)
    add_header(s8, "Safety & Compliance", "Enterprise Guardrails: 'No Means No'")

    guard_cards = [
        ("Anti-Bribery & Deceptive Marketing (Sec 13.6)", C_ACCENT_RED, [
            "Prohibits disguising government official entertainment under 'General Marketing' to avoid scrutiny.",
            "Refuses ticket creation even when user explicitly demands: 'create it anyway'.",
            "Mandates transparent Concur reporting + Manager Pre-Approval for client expenses > $100."
        ]),
        ("Data Loss Prevention (DLP) & Privacy", C_PRIMARY, [
            "Singapore NRIC Masking: In-flight sanitization replaces S1234567A with [NRIC_REDACTED].",
            "Credit Card & Token Scrubbing: Strips PANs and API keys before sending to LLM context.",
            "GDPR / Right to be Forgotten: Session deletion endpoint for full compliance."
        ]),
        ("Transaction & Leave Notice Guardrails", C_ACCENT_YELLOW, [
            "Overdraft Prevention: Blocks leave requests that exceed accrued balances.",
            "15-Day Advance Notice Guard: Enforces statutory notice for long vacations while allowing routine 1-2 day short leaves.",
            "Medical Certificate (MC) Invariant: Enforces MC submission within 48h for leaves > 2 days."
        ])
    ]

    left_g = Inches(0.8)
    for g_title, g_accent, g_bullets in guard_cards:
        add_card(s8, left_g, Inches(1.6), Inches(3.64), Inches(5.2))
        tb = s8.shapes.add_textbox(left_g + Inches(0.2), Inches(1.8), Inches(3.24), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = g_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = g_accent
        p.space_after = Pt(12)

        for b in g_bullets:
            pb = tf.add_paragraph()
            pb.text = f"• {b}"
            pb.font.size = Pt(11)
            pb.font.color.rgb = C_TEXT_MAIN
            pb.space_after = Pt(8)

        left_g += Inches(3.95)


    # =========================================================================
    #     # =========================================================================
    # SLIDE 9: Comprehensive 4-Tier Evaluation Harness
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, C_LIGHT_BG)
    add_header(s9, "Verification & Testing", "Comprehensive 4-Tier Evaluation Suite (33 Test Cases)")

    add_card(s9, Inches(0.8), Inches(1.6), Inches(4.0), Inches(5.2), bg_color=C_PRIMARY_LIGHT, border_color=C_PRIMARY)
    tb_sc = s9.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(3.6), Inches(4.7))
    tf_sc = tb_sc.text_frame
    tf_sc.word_wrap = True

    p = tf_sc.paragraphs[0]
    p.text = "Overall Readiness"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p_big = tf_sc.add_paragraph()
    p_big.text = "99.34%"
    p_big.font.size = Pt(36)
    p_big.font.bold = True
    p_big.font.color.rgb = C_ACCENT_GREEN
    p_big.space_after = Pt(8)

    p_status = tf_sc.add_paragraph()
    p_status.text = "✅ 33 / 33 Passed (100.0% Pass Rate)"
    p_status.font.size = Pt(12)
    p_status.font.bold = True
    p_status.font.color.rgb = C_TEXT_MAIN
    p_status.space_after = Pt(14)

    metrics = [
        ("Relevance (S_rel, 30%)", "99.8%"),
        ("Rigor & Reasoning (S_rigor, 35%)", "100.0%"),
        ("Cost & Time (S_cost_time, 15%)", "96.0%"),
        ("Safety Guardrails (S_guard, 20%)", "100.0%")
    ]

    for m_label, m_val in metrics:
        pm = tf_sc.add_paragraph()
        pm.text = f"• {m_label}: {m_val}"
        pm.font.size = Pt(11)
        pm.font.color.rgb = C_TEXT_MAIN
        pm.space_after = Pt(4)

    add_card(s9, Inches(5.1), Inches(1.6), Inches(7.4), Inches(5.2))
    tb_tiers = s9.shapes.add_textbox(Inches(5.35), Inches(1.85), Inches(6.9), Inches(4.7))
    tf_tiers = tb_tiers.text_frame
    tf_tiers.word_wrap = True

    p = tf_tiers.paragraphs[0]
    p.text = "Stratified Dataset Coverage"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.space_after = Pt(10)

    tier_breakdown = [
        ("🟢 Tier 1: Happy Path (10 Cases)", "MOM statutory sick/maternity leaves, vacation accrual calculations, WorkWeek balance lookups, and ServiceImmediately ticket submissions."),
        ("🟡 Tier 2: MAS Gotchas & Multi-Hop (6 Cases)", "Cross-specialist handoffs, exhaustion preconditions, prohibited venue overrides, and ticket priority anti-inflation."),
        ("🔴 Tier 3: Hallucination Baits (3 Cases)", "Fictitious perks (pet helicopter transfers, crypto dining stipends, corporate yacht charters) testing strict zero-hallucination abstention."),
        ("🟣 Tier 4: Boundary & Safety Probes (3 Cases)", "Out-of-scope domain probes (code algorithms, political opinions, stock trading advice) verifying polite refusals."),
        ("💬 Multi-Turn Trajectories (3 Flows)", "Complex stateful conversational flows with clarification loops and iterative issue troubleshooting."),
        ("🛡️ Adversarial & Guardrails (8 Cases)", "Singapore NRIC masking, credit card DLP, prompt injection jailbreak defenses, and SaaS 500 error escalation.")
    ]

    for t_name, t_desc in tier_breakdown:
        pt = tf_tiers.add_paragraph()
        pt.text = t_name
        pt.font.size = Pt(11.5)
        pt.font.bold = True
        pt.font.color.rgb = C_TEXT_MAIN

        pd = tf_tiers.add_paragraph()
        pd.text = t_desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = C_TEXT_MUTED
        pd.space_after = Pt(4)

    # =========================================================================
    # SLIDE 10: Production Deployment & Cloud Run Topology
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, C_LIGHT_BG)
    add_header(s10, "DevOps & Infrastructure", "Production Cloud Run Deployment Architecture")

    infra_cards = [
        ("Google Cloud Run", C_PRIMARY, [
            "Serverless auto-scaling container running FastAPI + Uvicorn.",
            "W3C Distributed Tracing headers propagated across all agent steps.",
            "Zero cold-start delay with optimized multi-threaded async loop."
        ]),
        ("Google Cloud Secret Manager", C_ACCENT_GREEN, [
            "Secure storage of FastMCP credentials (`hr-agent-mcp-token`).",
            "Dynamic binding during Cloud Run deployment (`--set-secrets=...`).",
            "Zero hardcoded secrets in source repository or container image."
        ]),
        ("Vertex AI & IAM Service Accounts", C_ACCENT_YELLOW, [
            "Direct IAM integration via Compute Engine Service Account.",
            "Zero credential expiration (eliminates local ADC login requirements).",
            "Enterprise Gemini 2.5 Flash LLM with low latency and high concurrency."
        ]),
        ("Automated CI/CD Deployment", C_TEXT_MAIN, [
            "1-Click deployment script: `./deploy_full_gcp.sh`.",
            "Supports automated secret versioning, rolling revisions, and dry-run preflight checks."
        ])
    ]

    pos_infra = [
        (Inches(0.8), Inches(1.6)),
        (Inches(6.8), Inches(1.6)),
        (Inches(0.8), Inches(4.35)),
        (Inches(6.8), Inches(4.35))
    ]

    for idx, (i_title, i_color, i_bullets) in enumerate(infra_cards):
        pos_l, pos_t = pos_infra[idx]
        add_card(s10, pos_l, pos_t, Inches(5.6), Inches(2.55))

        tb = s10.shapes.add_textbox(pos_l + Inches(0.2), pos_t + Inches(0.15), Inches(5.2), Inches(2.25))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = i_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = i_color
        p.space_after = Pt(8)

        for b in i_bullets:
            pb = tf.add_paragraph()
            pb.text = f"• {b}"
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = C_TEXT_MAIN
            pb.space_after = Pt(4)

    # =========================================================================
    # SLIDE 11: Summary & Business ROI
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11, C_DARK_BG)

    tb_end = s11.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(5.2))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True

    p0 = tf_end.paragraphs[0]
    p0.text = "ELEVATE HR • TEAM 12"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = C_ACCENT_YELLOW
    p0.space_after = Pt(8)

    p1 = tf_end.add_paragraph()
    p1.text = "Delivering True Enterprise Autonomous Value"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.space_after = Pt(20)

    roi_items = [
        ("Autonomous Self-Service", "Employees complete leave bookings, policy inquiries, and IT tickets in seconds without HR human intervention."),
        ("Multi-System Interoperability", "Seamless FastMCP connectivity bridges WorkWeek and ServiceImmediately into one conversational brain."),
        ("Zero-Compromise Security", "Immutable compliance invariants, Singapore NRIC masking, and Model Armor prompt injection protection."),
        ("Production Proven", "99.34% composite score across 33 rigorous benchmarks with full Google Cloud Run serverless scale.")
    ]

    for rt, rd in roi_items:
        pr = tf_end.add_paragraph()
        pr.text = f"🚀 {rt}: {rd}"
        pr.font.size = Pt(13.5)
        pr.font.color.rgb = RGBColor(226, 232, 240)
        pr.space_after = Pt(10)

    # Save Presentation
    prs.save(OUTPUT_PPTX)
    print(f"✅ Presentation successfully generated: {OUTPUT_PPTX}")


if __name__ == "__main__":
    create_presentation()
